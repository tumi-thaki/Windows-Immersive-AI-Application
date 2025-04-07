import pptx
from .baseExtractor import TextExtractor

class PowerPointExtractor(TextExtractor):
    def extract_text(self) -> str:
        prs = pptx.Presentation(self.file_path)
        text = ""
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text += shape.text + "\n"
        return self.clean_text(text)